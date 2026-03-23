#!/usr/bin/env python3
"""Generate mortly business proposal PowerPoint presentation (Korean)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──────────────────────────────────────────────
FOREST_800 = RGBColor(0x0F, 0x17, 0x29)
FOREST_700 = RGBColor(0x1F, 0x2D, 0x52)
FOREST_600 = RGBColor(0x2E, 0x3D, 0x68)
CREAM_100 = RGBColor(0xF8, 0xF7, 0xF4)
CREAM_300 = RGBColor(0xE5, 0xE2, 0xDC)
AMBER_400 = RGBColor(0xD4, 0xA8, 0x53)
AMBER_300 = RGBColor(0xE6, 0xC9, 0x6E)
AMBER_600 = RGBColor(0xA8, 0x81, 0x2E)
SAGE_500 = RGBColor(0x57, 0x62, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GOLD_BG = RGBColor(0xFD, 0xF9, 0xEF)
GREEN_CHECK = RGBColor(0x22, 0x8B, 0x22)
RED_X = RGBColor(0xCC, 0x33, 0x33)

# ── Fonts ─────────────────────────────────────────────────────
FONT_HEADING = "Calibri"
FONT_BODY = "Calibri"

# ── Paths ─────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
LOGO_PATH = os.path.join(PROJECT_DIR, "public", "logo", "logo.png")
OUTPUT_PATH = os.path.join(PROJECT_DIR, "mortly_business_proposal.pptx")

# ── Presentation setup ────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helper functions ──────────────────────────────────────────
def add_slide():
    """Add a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name=FONT_BODY, line_spacing=1.2):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline(slide, lines, left, top, width, height, font_size=16,
                  color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=FONT_BODY, line_spacing=1.5, bullet=False):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ("• " + line) if bullet else line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_gold_line(slide, left, top, width):
    """Add a gold accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = AMBER_400
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, bg_color, border_color=None):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_badge(slide, text, left, top, width=None, bg_color=AMBER_400,
              text_color=FOREST_800, font_size=10):
    """Add a small badge/pill."""
    w = width or Inches(1.5)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Pt(24)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_flow_step(slide, number, title, x, y, w=Inches(1.2), h=Inches(0.5),
                  color=AMBER_400, text_color=FOREST_800):
    """Add a flow step with number circle and title."""
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, Inches(0.45), Inches(0.45)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text(slide, title, x + Inches(0.55), y + Inches(0.05),
             w, h, font_size=13, color=FOREST_800, bold=False,
             font_name=FONT_BODY)


def add_arrow_right(slide, x, y, width=Inches(0.6)):
    """Add a right arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, x, y, width, Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = AMBER_400
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, x, y, height=Inches(0.3)):
    """Add a down arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.25), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = AMBER_400
    shape.line.fill.background()
    return shape


def add_logo(slide, left, top, height=Inches(0.6)):
    """Add logo if available."""
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, left, top, height=height)


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ══════════════════════════════════════════════════════════════
def slide_01_title():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    # Logo
    add_logo(slide, Inches(5.9), Inches(1.5), Inches(1.0))

    # Gold accent line
    add_gold_line(slide, Inches(5.5), Inches(2.8), Inches(2.3))

    # Tagline
    add_text(slide, "브로커들이 먼저 제안하는\n새로운 모기지 방식",
             Inches(1.5), Inches(3.1), Inches(10.3), Inches(1.5),
             font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    # Subtitle
    add_text(slide, "캐나다 최초 프라이버시 중심 모기지 마켓플레이스",
             Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.6),
             font_size=18, color=CREAM_300, bold=False,
             alignment=PP_ALIGN.CENTER)

    # Confidential
    add_text(slide, "Confidential  |  2026",
             Inches(4.5), Inches(6.6), Inches(4.3), Inches(0.4),
             font_size=11, color=SAGE_500, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 2 — 문제점
# ══════════════════════════════════════════════════════════════
def slide_02_problem():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "PROBLEM", Inches(0.8), Inches(0.6), Inches(1.2))
    add_text(slide, "모기지 브로커 탐색의 문제점",
             Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             font_size=32, color=FOREST_800, bold=True,
             font_name=FONT_HEADING)

    problems = [
        ("01", "신청인은 정보가 없다",
         "모기지 브로커를 찾으려면 은행에 전화하거나, 지인 소개에 의존하거나, 광고를 믿어야 합니다.\n투명한 비교 수단이 없습니다."),
        ("02", "브로커는 비효율적 영업에 의존",
         "라이선스를 가진 전문 브로커들이 고객 확보를 위해 과도한 시간과 비용을\n광고와 콜드콜에 투자하고 있습니다."),
        ("03", "개인정보가 노출된다",
         "브로커가 적합한지 알기도 전에 개인 금융 정보를 먼저 공유해야 합니다.\n프라이버시가 보장되지 않습니다."),
    ]

    for i, (num, title, desc) in enumerate(problems):
        y = Inches(2.2) + Inches(i * 1.6)
        # Number
        add_text(slide, num, Inches(0.8), y, Inches(0.6), Inches(0.5),
                 font_size=28, color=AMBER_400, bold=True, font_name=FONT_HEADING)
        # Title
        add_text(slide, title, Inches(1.6), y, Inches(10), Inches(0.5),
                 font_size=20, color=FOREST_800, bold=True)
        # Description
        add_text(slide, desc, Inches(1.6), y + Inches(0.45), Inches(10), Inches(0.8),
                 font_size=14, color=SAGE_500)


# ══════════════════════════════════════════════════════════════
# SLIDE 3 — 솔루션
# ══════════════════════════════════════════════════════════════
def slide_03_solution():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    add_badge(slide, "SOLUTION", Inches(5.5), Inches(1.0), Inches(1.4))

    add_text(slide, "mortly 소개",
             Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.8),
             font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    add_text(slide, "신청인이 익명으로 모기지 요청을 게시하면,\n인증된 브로커들이 경쟁적으로 제안하는 역방향 마켓플레이스",
             Inches(2), Inches(2.8), Inches(9.3), Inches(1.0),
             font_size=18, color=CREAM_300, alignment=PP_ALIGN.CENTER)

    add_gold_line(slide, Inches(5), Inches(4.0), Inches(3.3))

    # Three gold boxes
    labels = ["익명성", "경쟁 구도", "인증된 브로커"]
    for i, label in enumerate(labels):
        x = Inches(2.8) + Inches(i * 2.8)
        card = add_card(slide, x, Inches(4.5), Inches(2.2), Inches(0.9),
                        FOREST_700, AMBER_400)
        add_text(slide, label, x + Inches(0.1), Inches(4.65),
                 Inches(2.0), Inches(0.5),
                 font_size=18, color=AMBER_300, bold=True,
                 alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 4 — 이용 방법 (개요)
# ══════════════════════════════════════════════════════════════
def slide_04_how_it_works():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "HOW IT WORKS", Inches(0.8), Inches(0.6), Inches(1.6))
    add_text(slide, "간단한 3단계",
             Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             font_size=32, color=FOREST_800, bold=True, font_name=FONT_HEADING)
    add_text(slide, "모기지 브로커를 찾는 새로운 방법",
             Inches(0.8), Inches(1.9), Inches(11), Inches(0.5),
             font_size=16, color=SAGE_500)

    steps = [
        ("1", "요청서 제출", "나의 모기지 조건을 익명으로\n게시합니다. 개인정보 없이\n대출 금액, 부동산 유형, 타임라인\n등을 공유합니다."),
        ("2", "브로커 비교", "인증된 브로커들이 요청을 검토하고\n소개를 보냅니다. 프로필, 리뷰,\n전문 분야를 나란히 비교합니다."),
        ("3", "나의 조건에 맞게 연결", "적합한 브로커를 찾으면 연결합니다.\n내 신원은 내가 공개하기로\n결정할 때까지 비공개입니다."),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(1.0) + Inches(i * 4.0)

        # Card
        add_card(slide, x, Inches(2.8), Inches(3.4), Inches(3.8), WHITE, CREAM_300)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.3), Inches(3.1), Inches(0.7), Inches(0.7)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = AMBER_400
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(22)
        p.font.color.rgb = FOREST_800
        p.font.bold = True
        p.font.name = FONT_HEADING
        p.alignment = PP_ALIGN.CENTER

        # Title
        add_text(slide, title, x + Inches(0.2), Inches(4.0),
                 Inches(3.0), Inches(0.5),
                 font_size=18, color=FOREST_800, bold=True,
                 alignment=PP_ALIGN.CENTER)

        # Description
        add_text(slide, desc, x + Inches(0.3), Inches(4.6),
                 Inches(2.8), Inches(2.0),
                 font_size=13, color=SAGE_500,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.4)

        # Arrow between cards
        if i < 2:
            add_arrow_right(slide, x + Inches(3.5), Inches(4.5), Inches(0.4))


# ══════════════════════════════════════════════════════════════
# SLIDE 5 — 신청인 플로우
# ══════════════════════════════════════════════════════════════
def slide_05_borrower_flow():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "BORROWER FLOW", Inches(0.8), Inches(0.5), Inches(1.8))
    add_text(slide, "신청인 여정",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=30, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    steps = [
        "회원가입",
        "이메일 인증",
        "카테고리 선택\n(주거용 / 상업용)",
        "요청서 작성",
        "요청 공개 (익명)",
        "브로커 소개 수신",
        "브로커 비교",
        "선택 및 채팅",
        "평가 & 리뷰",
    ]

    # Horizontal flow in two rows
    for i, step in enumerate(steps):
        row = 0 if i < 5 else 1
        col = i if i < 5 else i - 5
        x = Inches(0.6) + Inches(col * 2.5)
        y = Inches(2.0) + Inches(row * 2.5)

        # Card
        add_card(slide, x, y, Inches(2.1), Inches(1.6), WHITE, CREAM_300)

        # Step number
        add_badge(slide, str(i + 1), x + Inches(0.05), y + Inches(0.1),
                  Inches(0.35), AMBER_400, FOREST_800, font_size=11)

        # Step text
        add_text(slide, step, x + Inches(0.1), y + Inches(0.5),
                 Inches(1.9), Inches(1.0),
                 font_size=12, color=FOREST_800, bold=False,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.3)

        # Arrow
        if i < 4:
            add_arrow_right(slide, x + Inches(2.15), y + Inches(0.6), Inches(0.3))
        elif i == 4:
            # Down arrow to second row
            add_arrow_down(slide, x + Inches(0.9), y + Inches(1.65), Inches(0.7))
        elif i < 8 and i > 4:
            add_arrow_right(slide, x + Inches(2.15), y + Inches(0.6), Inches(0.3))


# ══════════════════════════════════════════════════════════════
# SLIDE 6 — 브로커 플로우
# ══════════════════════════════════════════════════════════════
def slide_06_broker_flow():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "BROKER FLOW", Inches(0.8), Inches(0.5), Inches(1.6))
    add_text(slide, "브로커 여정",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=30, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    steps = [
        "회원가입",
        "온보딩\n(프로필 작성)",
        "관리자 인증\n대기",
        "요청 열람",
        "소개 전송\n(크레딧 1건 차감)",
        "신청인 선택\n대기",
        "채팅 시작",
        "평가 수신",
    ]

    for i, step in enumerate(steps):
        row = 0 if i < 4 else 1
        col = i if i < 4 else i - 4
        x = Inches(0.8) + Inches(col * 3.0)
        y = Inches(2.0) + Inches(row * 2.5)

        add_card(slide, x, y, Inches(2.5), Inches(1.6), WHITE, CREAM_300)

        add_badge(slide, str(i + 1), x + Inches(0.05), y + Inches(0.1),
                  Inches(0.35), AMBER_400, FOREST_800, font_size=11)

        add_text(slide, step, x + Inches(0.1), y + Inches(0.5),
                 Inches(2.3), Inches(1.0),
                 font_size=12, color=FOREST_800,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.3)

        if i < 3:
            add_arrow_right(slide, x + Inches(2.55), y + Inches(0.6), Inches(0.35))
        elif i == 3:
            add_arrow_down(slide, x + Inches(1.1), y + Inches(1.65), Inches(0.7))
        elif i < 7 and i > 3:
            add_arrow_right(slide, x + Inches(2.55), y + Inches(0.6), Inches(0.35))


# ══════════════════════════════════════════════════════════════
# SLIDE 7 — 관리자 플로우
# ══════════════════════════════════════════════════════════════
def slide_07_admin_flow():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "ADMIN FLOW", Inches(0.8), Inches(0.5), Inches(1.5))
    add_text(slide, "관리자 운영",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=30, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    items = [
        ("대시보드", "전체 현황 모니터링\n사용자, 요청, 대화 통계"),
        ("브로커 인증", "라이선스 확인 후\n승인 또는 거절"),
        ("요청 승인", "모기지 상담 요청\n검토 및 승인"),
        ("대화 감독", "활성 대화 모니터링\n부적절한 대화 관리"),
        ("신고 관리", "사용자 신고 접수\n조사 및 조치"),
        ("크레딧 관리", "브로커 크레딧\n수동 조정"),
        ("시스템 설정", "플랫폼 설정\n유지보수 모드"),
    ]

    for i, (title, desc) in enumerate(items):
        row = 0 if i < 4 else 1
        col = i if i < 4 else i - 4
        x = Inches(0.6) + Inches(col * 3.1)
        y = Inches(2.0) + Inches(row * 2.5)

        add_card(slide, x, y, Inches(2.7), Inches(1.8), WHITE, CREAM_300)
        add_text(slide, title, x + Inches(0.2), y + Inches(0.2),
                 Inches(2.3), Inches(0.4),
                 font_size=15, color=FOREST_800, bold=True,
                 alignment=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.2), y + Inches(0.7),
                 Inches(2.3), Inches(0.9),
                 font_size=11, color=SAGE_500,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    # Audit note
    add_text(slide, "※ 모든 관리자 행동은 감사 추적(Audit Trail)으로 기록됩니다",
             Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             font_size=12, color=AMBER_600, bold=True,
             alignment=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════
# SLIDE 8 — 플랫폼 스크린샷
# ══════════════════════════════════════════════════════════════
def slide_08_screenshots():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    add_text(slide, "mortly 플랫폼 경험",
             Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.8),
             font_size=32, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    add_text(slide, "모던하고 직관적인 사용자 인터페이스",
             Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.5),
             font_size=16, color=CREAM_300, alignment=PP_ALIGN.CENTER)

    # Placeholder mockup frames
    labels = ["홈페이지", "요청서 작성", "브로커 비교"]
    for i, label in enumerate(labels):
        x = Inches(1.0) + Inches(i * 3.9)
        # Frame
        frame = add_card(slide, x, Inches(2.8), Inches(3.5), Inches(3.6),
                         FOREST_700, SAGE_500)
        # Label
        add_text(slide, f"[ {label} 스크린샷 ]", x + Inches(0.3), Inches(4.3),
                 Inches(2.9), Inches(0.5),
                 font_size=14, color=SAGE_500, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 9 — 핵심 차별화 요소
# ══════════════════════════════════════════════════════════════
def slide_09_differentiators():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "DIFFERENTIATORS", Inches(0.8), Inches(0.5), Inches(1.9))
    add_text(slide, "mortly가 다른 이유",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=32, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    items = [
        ("프라이버시 중심",
         "신청인의 신원은 본인이 공개를 결정할 때까지 완전히 숨겨집니다. "
         "정확한 금액 대신 범위로 입력하여 프라이버시를 보호합니다.",
         "🔒"),
        ("역방향 마켓플레이스",
         "기존 방식과 반대로, 브로커가 먼저 신청인에게 제안합니다. "
         "소비자에게 선택권이 있는 구조입니다.",
         "🔄"),
        ("한영 이중 언어",
         "한국어와 영어를 완벽하게 지원합니다. "
         "캐나다 한인 커뮤니티를 위한 맞춤형 서비스입니다.",
         "🌏"),
        ("인증된 브로커만 참여",
         "모든 브로커는 관리자가 라이선스를 직접 확인합니다. "
         "미인증 브로커는 요청을 볼 수 없습니다.",
         "✓"),
    ]

    for i, (title, desc, icon) in enumerate(items):
        row = i // 2
        col = i % 2
        x = Inches(0.8) + Inches(col * 6.0)
        y = Inches(2.0) + Inches(row * 2.5)

        add_card(slide, x, y, Inches(5.5), Inches(2.1), WHITE, CREAM_300)

        # Icon
        add_text(slide, icon, x + Inches(0.3), y + Inches(0.3),
                 Inches(0.5), Inches(0.5), font_size=28, color=FOREST_800)

        # Title
        add_text(slide, title, x + Inches(0.9), y + Inches(0.3),
                 Inches(4.3), Inches(0.4),
                 font_size=18, color=FOREST_800, bold=True)

        # Desc
        add_text(slide, desc, x + Inches(0.9), y + Inches(0.85),
                 Inches(4.3), Inches(1.0),
                 font_size=13, color=SAGE_500, line_spacing=1.4)


# ══════════════════════════════════════════════════════════════
# SLIDE 10 — 시장 기회
# ══════════════════════════════════════════════════════════════
def slide_10_market():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    add_badge(slide, "MARKET", Inches(5.5), Inches(0.6), Inches(1.2))
    add_text(slide, "시장 기회",
             Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
             font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    stats = [
        ("~$2조", "캐나다 모기지 시장 규모"),
        ("~15,000명", "라이선스 보유 모기지 브로커"),
        ("35%", "브로커 경유 모기지 비율"),
        ("25만+", "재캐나다 한인 인구"),
        ("0", "역방향 마켓플레이스\n직접 경쟁자"),
    ]

    for i, (num, label) in enumerate(stats):
        x = Inches(0.5) + Inches(i * 2.5)
        y = Inches(2.8)

        add_card(slide, x, y, Inches(2.2), Inches(2.8), FOREST_700, FOREST_600)

        add_text(slide, num, x + Inches(0.1), y + Inches(0.4),
                 Inches(2.0), Inches(0.8),
                 font_size=32, color=AMBER_300, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

        add_text(slide, label, x + Inches(0.1), y + Inches(1.4),
                 Inches(2.0), Inches(1.0),
                 font_size=13, color=CREAM_300,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.3)


# ══════════════════════════════════════════════════════════════
# SLIDE 11 — 타겟 고객층
# ══════════════════════════════════════════════════════════════
def slide_11_segments():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "TARGET", Inches(0.8), Inches(0.5), Inches(1.2))
    add_text(slide, "타겟 고객층",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=32, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    # Borrower column
    add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.8), WHITE, CREAM_300)
    add_text(slide, "모기지 신청인", Inches(1.2), Inches(2.3),
             Inches(4.7), Inches(0.5),
             font_size=20, color=FOREST_800, bold=True)
    add_gold_line(slide, Inches(1.2), Inches(2.9), Inches(2.0))

    borrowers = [
        "첫 주택 구매자 — 브로커 선택이 막막한 분",
        "리파이낸싱 — 더 나은 조건을 찾는 분",
        "한인 커뮤니티 — 이중 언어 서비스가 필요한 분",
        "프라이버시 중시 — 개인정보 보호를 원하는 분",
        "상업용 부동산 — 사업 목적 대출이 필요한 분",
    ]
    add_multiline(slide, borrowers, Inches(1.2), Inches(3.2),
                  Inches(4.7), Inches(3.0),
                  font_size=14, color=FOREST_800, bullet=True, line_spacing=1.6)

    # Broker column
    add_card(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.8), WHITE, CREAM_300)
    add_text(slide, "모기지 브로커", Inches(7.4), Inches(2.3),
             Inches(4.7), Inches(0.5),
             font_size=20, color=FOREST_800, bold=True)
    add_gold_line(slide, Inches(7.4), Inches(2.9), Inches(2.0))

    brokers = [
        "개인 브로커 — 효율적 리드 확보가 필요한 분",
        "성장 중인 중개사 — 디지털 고객 확보를 원하는 곳",
        "한인 시장 전문 — 한인 고객층 특화 브로커",
        "신규 브로커 — 고객 기반을 구축하려는 분",
        "전문 분야 브로커 — 실력으로 승부하고 싶은 분",
    ]
    add_multiline(slide, brokers, Inches(7.4), Inches(3.2),
                  Inches(4.7), Inches(3.0),
                  font_size=14, color=FOREST_800, bullet=True, line_spacing=1.6)


# ══════════════════════════════════════════════════════════════
# SLIDE 12 — 수익 모델 개요
# ══════════════════════════════════════════════════════════════
def slide_12_revenue_overview():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    add_badge(slide, "REVENUE", Inches(5.5), Inches(0.6), Inches(1.3))
    add_text(slide, "수익 모델",
             Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
             font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    add_text(slide, "신청인은 항상 무료  ·  브로커 SaaS 과금 모델",
             Inches(2), Inches(2.2), Inches(9.3), Inches(0.6),
             font_size=18, color=CREAM_300, alignment=PP_ALIGN.CENTER)

    add_gold_line(slide, Inches(5), Inches(3.0), Inches(3.3))

    # Two revenue stream cards
    streams = [
        ("월 구독료", "FREE / BASIC / PRO / PREMIUM\n4단계 구독 요금제로\n지속적인 반복 수익 창출",
         "$29 ~ $129/월"),
        ("크레딧 팩", "월 크레딧 소진 시\n추가 크레딧 팩 구매 가능\n(BASIC/PRO 티어 대상)",
         "$29 ~ $79/팩"),
    ]

    for i, (title, desc, price) in enumerate(streams):
        x = Inches(1.5) + Inches(i * 5.5)
        add_card(slide, x, Inches(3.5), Inches(4.8), Inches(3.2),
                 FOREST_700, AMBER_400)

        add_text(slide, title, x + Inches(0.3), Inches(3.8),
                 Inches(4.2), Inches(0.5),
                 font_size=22, color=AMBER_300, bold=True,
                 alignment=PP_ALIGN.CENTER)

        add_text(slide, price, x + Inches(0.3), Inches(4.4),
                 Inches(4.2), Inches(0.4),
                 font_size=14, color=CREAM_300,
                 alignment=PP_ALIGN.CENTER)

        add_text(slide, desc, x + Inches(0.3), Inches(5.0),
                 Inches(4.2), Inches(1.5),
                 font_size=14, color=CREAM_300,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.4)


# ══════════════════════════════════════════════════════════════
# SLIDE 13 — 구독 요금제
# ══════════════════════════════════════════════════════════════
def slide_13_pricing():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "PRICING", Inches(0.8), Inches(0.5), Inches(1.2))
    add_text(slide, "브로커 구독 요금제",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=30, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    tiers = [
        {
            "name": "무료", "price": "$0", "original": None, "discount": None,
            "period": "", "features": ["요청 열람만 가능", "소개 전송 불가", "상담 가능 건수: 없음"],
            "highlight": False,
        },
        {
            "name": "베이직", "price": "$29", "original": "$49", "discount": "41% OFF",
            "period": "/월", "features": ["상담 가능 건수: 월 5건", "크레딧 팩 구매 가능"],
            "highlight": False,
        },
        {
            "name": "프로", "price": "$69", "original": "$99", "discount": "30% OFF",
            "period": "/월", "features": ["상담 가능 건수: 월 20건", "신규 상담요청 알림기능", "크레딧 팩 구매 가능"],
            "highlight": True,
        },
        {
            "name": "프리미엄", "price": "$129", "original": "$199", "discount": "35% OFF",
            "period": "/월", "features": ["상담 가능 건수: 무제한", "신규 상담요청 알림기능", "실시간 메시지 알림"],
            "highlight": False,
        },
    ]

    for i, tier in enumerate(tiers):
        x = Inches(0.5) + Inches(i * 3.15)
        y = Inches(1.9)
        w = Inches(2.9)
        h = Inches(5.0)

        border = AMBER_400 if tier["highlight"] else CREAM_300
        bg = FOREST_800 if tier["highlight"] else WHITE
        text_color = WHITE if tier["highlight"] else FOREST_800
        sub_color = CREAM_300 if tier["highlight"] else SAGE_500
        feat_color = CREAM_300 if tier["highlight"] else FOREST_800

        add_card(slide, x, y, w, h, bg, border)

        if tier["highlight"]:
            add_badge(slide, "가장 인기", x + Inches(0.7), y - Inches(0.15),
                      Inches(1.5), AMBER_400, FOREST_800, font_size=10)

        # Tier name
        add_text(slide, tier["name"], x + Inches(0.2), y + Inches(0.3),
                 Inches(2.5), Inches(0.4),
                 font_size=14, color=AMBER_400 if tier["highlight"] else SAGE_500,
                 bold=True)

        # Original price + discount
        if tier["original"]:
            add_text(slide, tier["original"], x + Inches(0.2), y + Inches(0.8),
                     Inches(0.8), Inches(0.3),
                     font_size=14, color=sub_color)
            # Strikethrough simulated with a line
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x + Inches(0.2), y + Inches(0.93),
                Inches(0.6), Pt(1.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = sub_color
            shape.line.fill.background()

            add_badge(slide, tier["discount"],
                      x + Inches(1.0), y + Inches(0.78),
                      Inches(1.1), RGBColor(0xE5, 0x3E, 0x3E), WHITE, font_size=9)

        # Price
        price_y = y + Inches(1.2) if tier["original"] else y + Inches(0.8)
        add_text(slide, tier["price"] + tier["period"],
                 x + Inches(0.2), price_y,
                 Inches(2.5), Inches(0.6),
                 font_size=32, color=text_color, bold=True,
                 font_name=FONT_HEADING)

        # Features
        feat_y = price_y + Inches(0.8)
        for j, feat in enumerate(tier["features"]):
            add_text(slide, "✓  " + feat,
                     x + Inches(0.25), feat_y + Inches(j * 0.4),
                     Inches(2.4), Inches(0.35),
                     font_size=12, color=feat_color)


# ══════════════════════════════════════════════════════════════
# SLIDE 14 — 크레딧 팩 및 단위 경제학
# ══════════════════════════════════════════════════════════════
def slide_14_credits():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "UNIT ECONOMICS", Inches(0.8), Inches(0.5), Inches(2.0))
    add_text(slide, "크레딧 팩 및 단위 경제학",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=30, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    # Credit packs
    packs = [
        ("SMALL 팩", "3 크레딧", "$29", "$9.67 / 건"),
        ("LARGE 팩", "10 크레딧", "$79", "$7.90 / 건"),
    ]

    for i, (name, credits, price, unit) in enumerate(packs):
        x = Inches(0.8) + Inches(i * 4.0)
        add_card(slide, x, Inches(2.0), Inches(3.5), Inches(2.2), WHITE, CREAM_300)

        add_text(slide, name, x + Inches(0.3), Inches(2.2),
                 Inches(2.9), Inches(0.4),
                 font_size=18, color=FOREST_800, bold=True)

        add_text(slide, f"{credits}  ·  {price}", x + Inches(0.3), Inches(2.7),
                 Inches(2.9), Inches(0.4),
                 font_size=16, color=SAGE_500)

        add_text(slide, unit, x + Inches(0.3), Inches(3.2),
                 Inches(2.9), Inches(0.4),
                 font_size=22, color=AMBER_400, bold=True, font_name=FONT_HEADING)

    # ROI box
    add_card(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.2),
             FOREST_800, AMBER_400)

    add_text(slide, "브로커 ROI 분석", Inches(1.2), Inches(4.9),
             Inches(10.5), Inches(0.5),
             font_size=20, color=AMBER_300, bold=True)

    roi_lines = [
        "캐나다 평균 모기지 브로커 수수료: $3,000 ~ $8,000 / 건",
        "mortly를 통한 고객 1명 획득 비용: $7.90 (LARGE 팩 기준)",
        "전환율 10% 가정 시, 실질 고객 획득 비용: ~$79",
        "ROI: 최소 38배 ~ 101배 투자 수익률",
    ]
    add_multiline(slide, roi_lines, Inches(1.2), Inches(5.5),
                  Inches(10.5), Inches(1.2),
                  font_size=14, color=CREAM_300, bullet=True, line_spacing=1.5)


# ══════════════════════════════════════════════════════════════
# SLIDE 15 — 매출 예상
# ══════════════════════════════════════════════════════════════
def slide_15_revenue_projections():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    add_badge(slide, "PROJECTIONS", Inches(5.2), Inches(0.6), Inches(1.6))
    add_text(slide, "매출 성장 시나리오",
             Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
             font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    projections = [
        ("1년차", "50명", "~$29K", Inches(1.5)),
        ("2년차", "200명", "~$142K", Inches(3.5)),
        ("3년차", "500명", "~$414K", Inches(5.5)),
    ]

    # Bar chart (simplified)
    bar_heights = [Inches(1.0), Inches(2.2), Inches(3.5)]
    bar_bottom = Inches(6.2)

    for i, ((year, brokers, arr, _), bh) in enumerate(zip(projections, bar_heights)):
        x = Inches(2.5) + Inches(i * 3.0)

        # Bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, bar_bottom - bh, Inches(2.0), bh
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = AMBER_400 if i == 2 else AMBER_300 if i == 1 else SAGE_500
        bar.line.fill.background()

        # ARR on bar
        add_text(slide, arr, x, bar_bottom - bh - Inches(0.5),
                 Inches(2.0), Inches(0.4),
                 font_size=24, color=AMBER_300, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

        # Year label
        add_text(slide, year, x, bar_bottom + Inches(0.1),
                 Inches(2.0), Inches(0.3),
                 font_size=14, color=CREAM_300, bold=True,
                 alignment=PP_ALIGN.CENTER)

        # Broker count
        add_text(slide, f"유료 브로커 {brokers}", x, bar_bottom + Inches(0.4),
                 Inches(2.0), Inches(0.3),
                 font_size=11, color=SAGE_500,
                 alignment=PP_ALIGN.CENTER)

    # Note
    add_text(slide, "* 크레딧 팩 매출 별도 (+약 20%)",
             Inches(1.5), Inches(7.0), Inches(10), Inches(0.3),
             font_size=11, color=SAGE_500, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 16 — 경쟁 환경
# ══════════════════════════════════════════════════════════════
def slide_16_competitive():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "COMPETITIVE", Inches(0.8), Inches(0.5), Inches(1.6))
    add_text(slide, "경쟁 환경 분석",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=30, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    # Table
    headers = ["기능", "mortly", "Ratehub", "RateShop", "Nesto", "기존 소개"]
    rows = [
        ["익명 요청", "✓", "✗", "✗", "✗", "✗"],
        ["역방향 마켓플레이스", "✓", "✗", "✗", "✗", "✗"],
        ["한영 이중 언어", "✓", "✗", "✗", "✗", "드물게"],
        ["브로커 경쟁 구도", "✓", "제한적", "제한적", "✗", "✗"],
        ["신청인 무료", "✓", "✓", "✓", "✓", "✓"],
        ["수동 브로커 인증", "✓", "다양", "다양", "자체", "없음"],
    ]

    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols,
        Inches(0.8), Inches(2.0),
        Inches(11.5), Inches(4.5)
    )
    table = table_shape.table

    # Style header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = FOREST_800

    # Style rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

            if j == 1:  # mortly column
                p.font.color.rgb = GREEN_CHECK if val == "✓" else FOREST_800
                p.font.bold = True
            elif val == "✓":
                p.font.color.rgb = GREEN_CHECK
            elif val == "✗":
                p.font.color.rgb = RED_X
            else:
                p.font.color.rgb = SAGE_500

            cell.fill.solid()
            cell.fill.fore_color.rgb = CREAM_100 if i % 2 == 0 else WHITE


# ══════════════════════════════════════════════════════════════
# SLIDE 17 — 기술 스택
# ══════════════════════════════════════════════════════════════
def slide_17_tech():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "TECHNOLOGY", Inches(0.8), Inches(0.5), Inches(1.6))
    add_text(slide, "확장 가능한 기술 기반",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=30, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    stack = [
        ("프론트엔드", "Next.js 16\nReact 19\nTypeScript\nTailwind CSS"),
        ("백엔드", "Next.js API Routes\nPrisma ORM 5\nNextAuth.js (JWT)"),
        ("데이터베이스", "PostgreSQL\nSupabase 호스팅\n실시간 WebSocket"),
        ("인프라", "Vercel 배포\n자동 스케일링\nEdge Network\nCron 자동화"),
        ("다국어", "next-i18next\n한국어 / 영어\n완전 이중 언어"),
    ]

    for i, (title, items) in enumerate(stack):
        x = Inches(0.5) + Inches(i * 2.5)
        add_card(slide, x, Inches(2.0), Inches(2.3), Inches(4.5), WHITE, CREAM_300)

        add_text(slide, title, x + Inches(0.2), Inches(2.2),
                 Inches(1.9), Inches(0.4),
                 font_size=15, color=AMBER_600, bold=True,
                 alignment=PP_ALIGN.CENTER)

        add_gold_line(slide, x + Inches(0.4), Inches(2.7), Inches(1.5))

        add_text(slide, items, x + Inches(0.2), Inches(3.0),
                 Inches(1.9), Inches(3.0),
                 font_size=13, color=FOREST_800,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.6)


# ══════════════════════════════════════════════════════════════
# SLIDE 18 — 플랫폼 현황
# ══════════════════════════════════════════════════════════════
def slide_18_status():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    add_badge(slide, "STATUS", Inches(5.5), Inches(0.6), Inches(1.2))
    add_text(slide, "플랫폼 현황: MVP 완성",
             Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
             font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    # Built column
    add_card(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.5),
             FOREST_700, FOREST_600)
    add_text(slide, "✓  구축 완료", Inches(1.2), Inches(2.6),
             Inches(4.7), Inches(0.5),
             font_size=18, color=AMBER_300, bold=True)

    built = [
        "신청인 전체 요청 플로우 (주거용 + 상업용)",
        "브로커 온보딩, 인증, 열람, 소개",
        "실시간 채팅 (Supabase Realtime)",
        "종합 관리자 패널 (대시보드, 관리, 감사)",
        "구독 티어 시스템 + 크레딧 관리",
        "이메일 인증 (Resend)",
        "한영 이중 언어 전체 지원",
        "모바일 반응형 디자인",
        "Vercel 배포 완료",
    ]
    add_multiline(slide, built, Inches(1.2), Inches(3.3),
                  Inches(4.7), Inches(3.5),
                  font_size=13, color=CREAM_300, bullet=True, line_spacing=1.4)

    # Next milestones column
    add_card(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.5),
             FOREST_700, FOREST_600)
    add_text(slide, "→  다음 마일스톤", Inches(7.4), Inches(2.6),
             Inches(4.7), Inches(0.5),
             font_size=18, color=AMBER_300, bold=True)

    nexts = [
        "Stripe 결제 연동",
        "푸시 알림 시스템",
        "OAuth 로그인 (Google)",
        "추천 브로커 배치 알고리즘",
        "문서 업로드 기능",
        "연간 결제 할인",
        "분석 대시보드 (Pro/Premium)",
    ]
    add_multiline(slide, nexts, Inches(7.4), Inches(3.3),
                  Inches(4.7), Inches(3.5),
                  font_size=13, color=CREAM_300, bullet=True, line_spacing=1.4)


# ══════════════════════════════════════════════════════════════
# SLIDE 19 — 신뢰 및 안전
# ══════════════════════════════════════════════════════════════
def slide_19_trust():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "TRUST & SAFETY", Inches(0.8), Inches(0.5), Inches(1.8))
    add_text(slide, "신뢰 및 안전 시스템",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=30, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    items = [
        ("수동 브로커 인증", "관리자가 모든 브로커의 라이선스를\n직접 확인 후 승인합니다"),
        ("감사 추적", "모든 관리자 행동이 타임스탬프,\n대상, 사유와 함께 기록됩니다"),
        ("신고 시스템", "사용자가 브로커나 대화를 신고하면\n관리자가 조사 후 조치합니다"),
        ("공개 ID 시스템", "9자리 숫자 ID로 식별하며\n내부 DB ID는 절대 노출하지 않습니다"),
        ("요청 관리", "관리자가 모든 요청을 검토 후\n승인, 거절, 종료할 수 있습니다"),
        ("대화 모니터링", "관리자가 활성 대화를 확인하고\n부적절한 대화를 종료할 수 있습니다"),
    ]

    for i, (title, desc) in enumerate(items):
        row = i // 3
        col = i % 3
        x = Inches(0.6) + Inches(col * 4.1)
        y = Inches(2.0) + Inches(row * 2.5)

        add_card(slide, x, y, Inches(3.7), Inches(2.1), WHITE, CREAM_300)
        add_text(slide, "🛡", x + Inches(0.2), y + Inches(0.2),
                 Inches(0.4), Inches(0.4), font_size=18, color=FOREST_800)
        add_text(slide, title, x + Inches(0.6), y + Inches(0.2),
                 Inches(2.8), Inches(0.4),
                 font_size=15, color=FOREST_800, bold=True)
        add_text(slide, desc, x + Inches(0.3), y + Inches(0.8),
                 Inches(3.1), Inches(1.0),
                 font_size=12, color=SAGE_500, line_spacing=1.4)


# ══════════════════════════════════════════════════════════════
# SLIDE 20 — 시장 진출 전략
# ══════════════════════════════════════════════════════════════
def slide_20_gtm():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    add_badge(slide, "GO-TO-MARKET", Inches(5.0), Inches(0.6), Inches(1.8))
    add_text(slide, "시장 진출 전략",
             Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.8),
             font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    phases = [
        ("1단계: 런칭", "GTA 한인 커뮤니티",
         "• 이중 언어 우위 활용\n• 한인 미디어 및 소셜 마케팅\n• 커뮤니티 파트너십\n• 초기 브로커 50명 확보"),
        ("2단계: 확장", "캐나다 전역 확대",
         "• 밴쿠버, 캘거리 진출\n• SEO 및 콘텐츠 마케팅\n• 브로커 추천 프로그램\n• 전체 캐나다 신청인 대상"),
        ("3단계: 스케일", "전국 + 파트너십",
         "• 전국 커버리지\n• 부동산 플랫폼 파트너십\n• 새로운 언어 커뮤니티\n  (중국어, 펀잡어 등)\n• 프리미엄 기능 고도화"),
    ]

    for i, (title, subtitle, items) in enumerate(phases):
        x = Inches(0.6) + Inches(i * 4.2)
        add_card(slide, x, Inches(2.4), Inches(3.8), Inches(4.5),
                 FOREST_700, FOREST_600)

        add_text(slide, title, x + Inches(0.3), Inches(2.6),
                 Inches(3.2), Inches(0.4),
                 font_size=18, color=AMBER_300, bold=True)

        add_text(slide, subtitle, x + Inches(0.3), Inches(3.1),
                 Inches(3.2), Inches(0.4),
                 font_size=14, color=WHITE, bold=True)

        add_gold_line(slide, x + Inches(0.3), Inches(3.6), Inches(2.5))

        add_text(slide, items, x + Inches(0.3), Inches(3.9),
                 Inches(3.2), Inches(2.8),
                 font_size=13, color=CREAM_300, line_spacing=1.5)

        if i < 2:
            add_arrow_right(slide, x + Inches(3.9), Inches(4.5), Inches(0.3))


# ══════════════════════════════════════════════════════════════
# SLIDE 21 — 팀 소개
# ══════════════════════════════════════════════════════════════
def slide_21_team():
    slide = add_slide()
    set_bg(slide, CREAM_100)

    add_badge(slide, "TEAM", Inches(0.8), Inches(0.5), Inches(1.0))
    add_text(slide, "팀 소개",
             Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
             font_size=32, color=FOREST_800, bold=True, font_name=FONT_HEADING)

    # Placeholder team cards
    for i in range(3):
        x = Inches(1.0) + Inches(i * 4.0)
        add_card(slide, x, Inches(2.2), Inches(3.5), Inches(4.5), WHITE, CREAM_300)

        # Photo placeholder
        photo = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.0), Inches(2.6),
            Inches(1.5), Inches(1.5)
        )
        photo.fill.solid()
        photo.fill.fore_color.rgb = CREAM_300
        photo.line.fill.background()

        add_text(slide, "👤", x + Inches(1.3), Inches(2.9),
                 Inches(1.0), Inches(0.8),
                 font_size=36, color=SAGE_500, alignment=PP_ALIGN.CENTER)

        add_text(slide, "[이름]", x + Inches(0.3), Inches(4.4),
                 Inches(2.9), Inches(0.4),
                 font_size=18, color=FOREST_800, bold=True,
                 alignment=PP_ALIGN.CENTER)

        add_text(slide, "[직함]", x + Inches(0.3), Inches(4.9),
                 Inches(2.9), Inches(0.3),
                 font_size=14, color=AMBER_600,
                 alignment=PP_ALIGN.CENTER)

        add_text(slide, "[약력 및 전문 분야]", x + Inches(0.3), Inches(5.3),
                 Inches(2.9), Inches(1.0),
                 font_size=12, color=SAGE_500,
                 alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 22 — 마무리 / 연락처
# ══════════════════════════════════════════════════════════════
def slide_22_closing():
    slide = add_slide()
    set_bg(slide, FOREST_800)

    # Logo
    add_logo(slide, Inches(5.9), Inches(1.8), Inches(1.0))

    add_gold_line(slide, Inches(5.5), Inches(3.1), Inches(2.3))

    add_text(slide, "브로커들이 먼저 제안하는\n새로운 모기지 방식",
             Inches(1.5), Inches(3.4), Inches(10.3), Inches(1.2),
             font_size=30, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name=FONT_HEADING)

    add_text(slide, "mortly.ca",
             Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
             font_size=18, color=AMBER_300, bold=True,
             alignment=PP_ALIGN.CENTER)

    add_text(slide, "support@mortly.ca",
             Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.5),
             font_size=14, color=CREAM_300,
             alignment=PP_ALIGN.CENTER)

    add_text(slide, "Confidential  ·  내부 검토용",
             Inches(4.5), Inches(6.6), Inches(4.3), Inches(0.4),
             font_size=11, color=SAGE_500, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# GENERATE
# ══════════════════════════════════════════════════════════════
def main():
    print("Generating mortly business proposal presentation...")

    slide_01_title()
    slide_02_problem()
    slide_03_solution()
    slide_04_how_it_works()
    slide_05_borrower_flow()
    slide_06_broker_flow()
    slide_07_admin_flow()
    slide_08_screenshots()
    slide_09_differentiators()
    slide_10_market()
    slide_11_segments()
    slide_12_revenue_overview()
    slide_13_pricing()
    slide_14_credits()
    slide_15_revenue_projections()
    slide_16_competitive()
    slide_17_tech()
    slide_18_status()
    slide_19_trust()
    slide_20_gtm()
    slide_21_team()
    slide_22_closing()

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
